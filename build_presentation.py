"""Deep-dive generator for the UC-2 / UC-15 presentation deck. Run with the venv python."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ───────────────────────── theme ─────────────────────────
NAVY      = RGBColor(0x10, 0x2A, 0x43)   # header text / accent — "blue"
NAVY_DARK = RGBColor(0x0A, 0x1C, 0x2E)   # cover / divider / closing background
TEAL      = RGBColor(0x17, 0xA2, 0x8B)
ORANGE    = RGBColor(0xE3, 0x7B, 0x1B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY      = RGBColor(0x8A, 0x93, 0x9B)
ROW_ALT   = RGBColor(0xE9, 0xED, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.5)
CONTENT_W = SLIDE_W - 2 * MARGIN

FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

slide_counter = {"n": 0}


def new_slide():
    return prs.slides.add_slide(BLANK)


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size=18, color=TEXT_DARK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_bullet_list_at(slide, left, top, width, height, items, size=15, color=TEXT_DARK,
                        space_after=8, marker="–"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"{marker}   {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
        p.space_after = Pt(space_after)
        p.line_spacing = 1.05
    return box


def add_footer(slide, label):
    slide_counter["n"] += 1
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Pt(1.2), GRAY)
    add_text(slide, MARGIN, SLIDE_H - Inches(0.32), Inches(8), Inches(0.32),
              label, size=10, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.32), Inches(0.8), Inches(0.32),
              str(slide_counter["n"]), size=10, color=GRAY, align=PP_ALIGN.RIGHT,
              anchor=MSO_ANCHOR.MIDDLE)


def add_title_bar(slide, title, accent=NAVY):
    """Clean white header zone with a thin colored accent underline and navy/blue title text."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), WHITE)
    add_rect(slide, 0, Inches(1.0), SLIDE_W, Pt(2.5), accent)
    add_text(slide, MARGIN, 0, CONTENT_W, Inches(1.0), title, size=26, bold=True,
              color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ───────────────────────── slide builders ─────────────────────────

def title_slide(title, subtitle, tag, footer_text):
    slide = new_slide()
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_rect(slide, 0, Inches(4.55), SLIDE_W, Inches(0.06), TEAL)
    add_text(slide, Inches(1), Inches(2.4), Inches(11.33), Inches(1.3), title,
              size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1), Inches(3.6), Inches(11.33), Inches(0.7), subtitle,
              size=19, color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1), Inches(4.85), Inches(11.33), Inches(0.5), tag,
              size=14, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1), Inches(6.7), Inches(11.33), Inches(0.4), footer_text,
              size=12, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


def divider_slide(part_label, title, accent=NAVY):
    slide = new_slide()
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_rect(slide, Inches(1), Inches(3.55), Inches(2.2), Inches(0.07), accent)
    add_text(slide, Inches(1), Inches(2.6), Inches(11.33), Inches(0.6), part_label,
              size=18, bold=True, color=accent)
    add_text(slide, Inches(1), Inches(3.75), Inches(11.33), Inches(1.2), title,
              size=32, bold=True, color=WHITE)
    add_footer(slide, "Regulatory Compliance Automation — UC-2 & UC-15")
    return slide


def bullets_slide(title, bullets, accent=NAVY, notes=None, size=16.5, space_after=14):
    slide = new_slide()
    add_title_bar(slide, title, accent)
    box = slide.shapes.add_textbox(MARGIN, Inches(1.4), CONTENT_W, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"•   {item.strip()}"
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT_DARK
        run.font.name = FONT
        p.space_after = Pt(space_after)
        p.line_spacing = 1.1
    add_footer(slide, "Regulatory Compliance Automation — UC-2 & UC-15")
    if notes:
        set_notes(slide, notes)
    return slide


def agenda_slide():
    slide = new_slide()
    add_title_bar(slide, "Agenda", NAVY)

    groups = [
        ("UC-2", TEAL, [
            "Problem Statement",
            "Crawler Design: One Tool Per Site",
            "Extraction Strategy: Local vs. Cloud, Routed by Need",
            "System Architecture, API, and Scheduling",
            "Real Usage at Scale and Lessons Learned",
        ]),
        ("UC-15", ORANGE, [
            "Problem Statement and Model Evaluation",
            "The Three-Stage Extraction Pipeline",
            "Hallucination Control and Reliability",
            "Matching, Gap Analysis, and Lessons Learned",
        ]),
        ("Costing & Wrap-Up", NAVY, [
            "Costing — Per-Document Economics",
            "What Isn't Instrumented",
            "Licensing, Infrastructure, and Key Takeaways",
        ]),
    ]

    y = Inches(1.3)
    for heading, color, items in groups:
        add_text(slide, MARGIN, y, Inches(4), Inches(0.36), heading, size=18, bold=True, color=color)
        y += Inches(0.4)
        item_h = Inches(0.32) * len(items) + Inches(0.1)
        add_bullet_list_at(slide, MARGIN + Inches(0.35), y, CONTENT_W - Inches(0.35), item_h,
                            items, size=14.5, space_after=5)
        y += item_h + Inches(0.18)

    add_footer(slide, "Regulatory Compliance Automation — UC-2 & UC-15")
    set_notes(slide, "UC-2 is the outer pipeline; UC-15 is the AI engine embedded inside its analysis step.")
    return slide


def table_slide(title, headers, rows, col_ratios, accent=NAVY, notes=None,
                 header_size=14, body_size=13, caption=None,
                 max_row_height=None, extra_bullets=None):
    slide = new_slide()
    add_title_bar(slide, title, accent)
    top = Inches(1.3)
    avail_h = SLIDE_H - top - Inches(0.55)

    n_rows = len(rows) + 1
    if max_row_height is not None:
        table_h = min(avail_h, Emu(int(max_row_height) * n_rows))
    else:
        reserve = Inches(0.5) if caption else 0
        if extra_bullets:
            reserve += Inches(1.5)
        table_h = avail_h - reserve

    gshape = slide.shapes.add_table(n_rows, len(headers), MARGIN, top, CONTENT_W, table_h)
    table = gshape.table

    total_units = sum(col_ratios)
    for i, ratio in enumerate(col_ratios):
        table.columns[i].width = Emu(int(CONTENT_W * ratio / total_units))

    row_h = Emu(int(table_h / n_rows))
    for r in range(n_rows):
        table.rows[r].height = row_h

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.07)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = FONT

    for r, row_vals in enumerate(rows, start=1):
        fill = WHITE if r % 2 == 1 else ROW_ALT
        for c, val in enumerate(row_vals):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = TEXT_DARK
                    run.font.name = FONT
                    run.font.bold = (c == 0)

    cursor = top + table_h + Inches(0.08)
    if caption:
        add_text(slide, MARGIN, cursor, CONTENT_W, Inches(0.45),
                  caption, size=11, italic=True, color=GRAY)
        cursor += Inches(0.48)

    if extra_bullets:
        add_bullet_list_at(slide, MARGIN, cursor, CONTENT_W, Inches(1.6),
                            extra_bullets, size=12.5, color=TEXT_DARK, space_after=6)

    add_footer(slide, "Regulatory Compliance Automation — UC-2 & UC-15")
    if notes:
        set_notes(slide, notes)
    return slide


def pipeline_slide(title, stages, accent=TEAL, caption=None, notes=None, sub_caption=None):
    slide = new_slide()
    add_title_bar(slide, title, accent)
    if sub_caption:
        add_text(slide, MARGIN, Inches(1.15), CONTENT_W, Inches(0.4), sub_caption,
                  size=12.5, italic=True, color=GRAY)

    n = len(stages)
    arrow_w = Inches(0.3)
    box_h = Inches(1.85)
    total_arrow = arrow_w * (n - 1)
    box_w = Emu(int((CONTENT_W - total_arrow) / n))
    y = Inches(2.85)

    x = MARGIN
    for i, (header, desc) in enumerate(stages):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = accent
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.07)
        tf.margin_right = Inches(0.07)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = f"{i+1}. {header}"
        r1.font.size = Pt(13.5)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r1.font.name = FONT
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT

        x += box_w
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y + box_h / 2 - Inches(0.17),
                                            arrow_w, Inches(0.34))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()
            arrow.shadow.inherit = False
            x += arrow_w

    if caption:
        add_text(slide, MARGIN, y + box_h + Inches(0.3), CONTENT_W, Inches(1.4),
                  caption, size=12.5, color=TEXT_DARK)

    add_footer(slide, "Regulatory Compliance Automation — UC-2 & UC-15")
    if notes:
        set_notes(slide, notes)
    return slide


# ───────────────────────── PART 0 — TITLE / AGENDA ─────────────────────────

title_slide(
    "Regulatory Compliance Automation",
    "UC-2: Compliance Pipeline and UC-15: Requirement Extraction Engine",
    "An architecture deep dive, covering design decisions and a usage-grounded cost model",
    "Prepared June 2026",
)

agenda_slide()

# ───────────────────────── PART 1 — UC-2 ─────────────────────────

divider_slide("PART 1", "UC-2 — Regulatory Compliance Pipeline", accent=TEAL)

bullets_slide(
    "UC-2 — Problem Statement",
    [
        "UC-2 was created to solve a documented business problem: regulatory monitoring was being performed manually by a team of ten staff, with frequent missed updates, misclassifications, and error-prone manual data entry.",
        "Maintaining this process across multiple regulators — originally SBP and SECP, now extended to include CBB and SAMA — had become time-consuming and difficult to sustain at the level of accuracy compliance requires.",
        "The target future state was fully automated detection of new and changed publications, structured and accurate data, complete auditability, and a significant reduction in cost.",
        "UC-2 delivers that future state: it crawls continuously, extracts content reliably, applies AI-driven analysis through UC-15, preserves a complete version history, and serves the results bilingually through one API.",
    ],
    accent=TEAL,
    notes="Business Problem and Future State framing pulled directly from the original UC-002 project document (December 2024) — the four-regulator scope reflects how the project has since grown beyond its original SBP/SECP starting point.",
)

table_slide(
    "UC-2 — Crawler Design: One Tool Per Site",
    ["Regulator", "What's Crawled", "Site Behavior That Forced the Choice", "Framework"],
    [
        ["CBB", "Rulebook modules, AML law, corporate governance, capital market regulations",
         "Mostly static, server-rendered HTML with a deep folder/chapter hierarchy",
         "Requests + BeautifulSoup, 2Captcha"],
        ["SAMA", "Circulars, laws, implementing regulations",
         "JS-rendered navigation and listing pages",
         "Selenium, 2Captcha"],
        ["SBP", "Circulars, notifications, regulatory returns, laws & regulations",
         "Large paginated archive — suits a structured spider/middleware model",
         "Scrapy, 2Captcha middleware"],
        ["SECP", "Acts, ordinances, rules, directives, circulars",
         "Modern JS-heavy site needing reliable headless rendering at scale",
         "Playwright, 2Captcha"],
    ],
    col_ratios=[1.0, 3.0, 3.6, 2.0],
    accent=TEAL,
    caption="Each crawler framework was matched to its regulator's actual site structure rather than forcing a single scraper across all four — a decision confirmed directly against the codebase.",
    notes="This wasn't an arbitrary tooling choice — each site's rendering behavior dictated the right framework.",
)

table_slide(
    "UC-2 — Extraction Strategy: Local vs. Cloud, Routed by Need",
    ["Dimension", "Original Design (Local-Only)", "Current Design (Routed by Need)"],
    [
        ["Formatting fidelity",
         "Tables, multi-column layouts and nested headers frequently collapsed or merged when extracted locally",
         "PDF.co (cloud, format-preserving) is primary wherever the output must stay structured HTML"],
        ["Cost",
         "Free — local compute only (PyMuPDF for digital, Tesseract OCR for scanned)",
         "PDF.co is paid per page/credit, but only spent where formatting genuinely matters"],
    ],
    col_ratios=[1.6, 3.8, 4.2],
    accent=TEAL,
    max_row_height=Inches(1.2),
    caption="This reflects real implementation history: the local-only approach was tried first and failed specifically on formatting fidelity, not on raw text accuracy — which is exactly why PDF.co became the primary engine for the formatting-sensitive path.",
    extra_bullets=[
        "Local extraction is still used today wherever the output only needs to feed the LLM as plain text — formatting doesn't matter there, so the free path is faster and PDF.co's per-page cost isn't justified.",
        "If PDF.co fails for any reason, the pipeline automatically falls back to local PyMuPDF for digital PDFs or Tesseract OCR for scanned PDFs, so a single vendor issue doesn't stop extraction.",
    ],
    notes="This is a genuine pivot, not a hypothetical comparison: local OCR was good enough for plain text but not for preserving table/column structure.",
)

pipeline_slide(
    "UC-2 — System Architecture",
    [
        ("Crawl", "Fetch documents per regulator"),
        ("Filter", "Deduplicate against existing records"),
        ("Insert\n(Metadata)", "Store the regulation and its metadata"),
        ("Extract", "Resolve full text via the 3-tier strategy"),
        ("Analyze", "UC-15: 3-stage LLM extraction"),
        ("Match", "Cross-reference existing controls and KPIs"),
    ],
    accent=TEAL,
    caption="Each phase is explained in detail on the next slide.",
    notes="One orchestrator class, one unified analysis table — the per-regulator differences are isolated to crawling and the insert/versioning step.",
)

bullets_slide(
    "UC-2 — System Architecture: Phase Details",
    [
        "Crawl: each regulator has its own crawler, covered earlier, that fetches documents from its website on a schedule and hands them to the orchestrator as a list of candidate documents.",
        "Filter: the orchestrator checks each candidate against what's already stored — by source URL for CBB, or by title, publication date, and document path for the others — so only new or modified documents continue.",
        "Insert with Metadata: a database record is created for the document, capturing its regulator, category, title, reference number, publication date, and source URL. For CBB specifically, this step also versions the content: because CBB's rulebook gets edited rather than only republished, any change first archives the prior content and prior analysis before writing the new version, so nothing is silently overwritten. The database currently holds 29,471 CBB content versions for 29,469 CBB documents, with only two marked inactive — meaning exactly one document has actually been revised since the initial crawl. SAMA, SBP, and SECP are not versioned this way; their documents are treated as point-in-time snapshots.",
        "Extract: the document's full text is resolved through a 3-tier strategy, so the analysis stage always receives usable text without an unnecessary live download. Tier 1 checks for text already extracted and cached, such as SAMA's pre-OCR'd PDF text or CBB's stored content. Tier 2 falls back to the stored HTML column. Tier 3, used only if the first two come up short, downloads the document live and extracts it, applying OCR if needed. Each tier only runs if the previous one returned fewer than two hundred characters, so most documents resolve at Tier 1 or 2.",
        "Analyze: the extracted text is handed to UC-15, the 3-stage LLM extraction engine. The model used is DeepSeek v3.2 via OpenRouter, chosen because it delivers balanced performance on classification and structuring tasks at a meaningfully lower cost than Claude or GPT-4-class models — important for a workload that calls the model multiple times per document across thousands of documents.",
        "Match: each extracted requirement is cross-referenced against the bank's existing requirements, controls, and KPIs, and is flagged as fully matched, partially matched, or new.",
    ],
    accent=TEAL,
    size=14,
    space_after=9,
)

bullets_slide(
    "UC-2 — API and Scheduling",
    [
        "The FastAPI application exposes a trigger endpoint for each regulator, along with endpoints for gap analysis and requirement mapping.",
        "A background thread checks a pipeline_schedule table every thirty seconds and starts the appropriate pipeline when it is due, using a single lock so that only one pipeline runs at a time.",
        "A second thread records a heartbeat every five minutes while a pipeline is running, so a stalled run becomes visible rather than failing silently.",
        "Scheduling is fully configurable per regulator through the API: a client can set a specific time of day, specific days of the week, a recurring interval such as every few hours, or multiple runs per day for any regulator, independently of the others.",
        "On-demand trigger endpoints let a client run any regulator's pipeline, or all of them, immediately — independent of its configured schedule.",
        "English and Arabic responses are cached per endpoint and automatically invalidated whenever the underlying analysis changes.",
        "Translation is handled through Google Translate, batched per response rather than per field, to keep the number of round trips low.",
    ],
    accent=TEAL,
    notes="The scheduler is intentionally simple — a polling loop with a lock, not a task queue — which fits the current single-pipeline-at-a-time operating model.",
)

table_slide(
    "UC-2 — Real Usage at Scale",
    ["Regulator", "Docs Crawled", "Docs LLM-Analyzed", "Note"],
    [
        ["CBB", "29,469", "9", "Pages below a folder depth of two are skipped — almost the entire crawl is navigation structure, not regulatory text"],
        ["SBP", "8,081", "509", "—"],
        ["SECP", "3,767", "195", "—"],
        ["SAMA", "675", "671", "Near one-to-one — almost every crawled SAMA document is analyzed"],
        ["Total", "41,992", "1,384", "—"],
    ],
    col_ratios=[1.0, 1.6, 1.8, 5.4],
    accent=TEAL,
    caption="These figures come directly from the regulations and compliance_analysis tables. CBB's crawl volume is mostly site structure, while SAMA, SBP, and SECP are much closer to a one-to-one ratio with analysis.",
    notes="The CBB gap (29,469 to 9) is the single most important real-usage fact in this deck for cost planning — almost none of that crawl volume reaches the LLM.",
)

bullets_slide(
    "UC-2 — Lessons Learned",
    [
        "Local OCR alone could not preserve formatting fidelity: multi-column tables and nested headers degraded enough that PDF.co became the primary method wherever the output needed to stay structured, such as the display HTML and versioned snapshots.",
        "Local extraction was not replaced, however — it remains the default for the path where the output only feeds the LLM as plain text, since formatting does not matter there and PDF.co's per-page cost is not justified.",
        "Each crawler framework was matched to its regulator's actual site behavior — static HTML, JavaScript rendering, or a paginated archive — rather than forcing one scraper across all four.",
        "CBB's depth-based filter, which skips analysis for pages above a certain folder depth, is why 29,469 crawled pages produced only nine analyses; without it, the pipeline would spend LLM calls on navigation and folder pages.",
    ],
    accent=TEAL,
    notes="These are the load-bearing decisions for cost and reliability — confirm with the team whether there's a similar story behind the unified compliance_analysis schema replacing per-regulator code paths.",
)

# ───────────────────────── PART 2 — UC-15 ─────────────────────────

divider_slide("PART 2", "UC-15 — Requirement Extraction Engine", accent=ORANGE)

bullets_slide(
    "UC-15 — Problem Statement",
    [
        "Regulatory text is written as prose — paragraphs, cross-references, and conditional clauses — and is not directly usable by a compliance system.",
        "UC-15's role is to convert that prose into atomic, testable obligations, classify each one, design a control for the obligations that require one, and cross-reference the result against the bank's existing requirements.",
        "The same engine processes every regulator's crawled text as well as manually uploaded documents, using a single code path rather than four separate ones.",
        "Output remains in the source language, English or Arabic, because UC-15 is a classification and structuring engine, not a translation engine.",
    ],
    accent=ORANGE,
    notes="UC-15 is the requirement-extraction part embedded inside UC-2's analysis layer — same engine regardless of where the text came from.",
)

table_slide(
    "UC-15 — Model Evaluation",
    ["Model", "Why It Wasn't Chosen", "Token Price\n(Input / Output per 1M)"],
    [
        ["Claude 3.5\nSonnet",
         "Strongest analytical depth of the five — the only model that extracted technical assessments as actionable requirements. Rejected on cost: at this price, running it several times per document across thousands of documents would cost roughly 10–37x more than the chosen model, for a quality gain this workload doesn't need.",
         "$3.00 / $15.00"],
        ["GPT-4.1",
         "Strong KPI and audit-trail output, but assigned departments too narrowly and missed technical depth. Still priced 7–20x higher than the chosen model.",
         "$2.00 / $8.00"],
        ["Gemini 2.5\nFlash",
         "Input pricing close to the chosen model, but output pricing is roughly 6x higher — and it missed technical assessments and specific regulatory consequences. The quality gap wasn't offset by a real cost advantage.",
         "$0.30 / $2.50"],
        ["Qwen 2.5\n72B",
         "Pricing is comparable to the chosen model, but it produced the weakest output of the five — the most generic controls, the weakest KPIs, the least specificity. Rejected on quality, not cost.",
         "$0.36 / $0.40"],
    ],
    col_ratios=[1.3, 6.2, 1.7],
    accent=ORANGE,
    header_size=12,
    body_size=11,
    max_row_height=Inches(1.05),
    caption="DeepSeek v3.2 — $0.27 / $0.40 per million tokens — was the model that didn't force a trade-off between quality and cost at this volume.",
    extra_bullets=[
        "Against the models near its price point (Gemini 2.5 Flash, Qwen 2.5 72B), DeepSeek produced the most accurate department and risk classification — the dimension that matters most for routing obligations to the right owner.",
        "Against the models with stronger qualitative depth (Claude 3.5 Sonnet, GPT-4.1), DeepSeek avoided a 7x to 37x cost premium that this workload — multiple calls per document, across thousands of documents — would make very expensive, very quickly.",
        "Net result: a model accurate enough for structured classification and control design, at a cost that scales comfortably with regulatory volume instead of against it.",
    ],
    notes="This is the team's actual side-by-side test across five models on real regulatory text. Pricing is OpenRouter's public per-token rate as of June 2026 — confirm current rates before using these figures for budgeting.",
)

pipeline_slide(
    "UC-15 — The Three-Stage Pipeline",
    [
        ("Extract", "Pull binding obligations such as must/shall/required; group into four to eight requirement clusters"),
        ("Normalize &\nClassify", "Dedupe; tag type, criticality, evidence, and execution category"),
        ("Control\nDesign", "Conditional — runs only for Ongoing Control obligations"),
    ],
    accent=ORANGE,
    caption="The model used throughout is DeepSeek v3.2 via OpenRouter, chosen for balanced classification performance at a meaningfully lower cost than Claude or GPT-4-class models. Stage 3 only runs when Stage 2 finds at least one Ongoing Control obligation, so not every document triggers a third call.",
    notes="processor/staged_LLM_Analyzer.py — described here as a 3-stage extraction engine (extract, classify, design controls); the pipeline also produces an internal executive-summary report from these three stages' output, not covered in this deck.",
)

bullets_slide(
    "UC-15 — Stage 1: Extraction Rules",
    [
        "Stage 1 extracts only binding obligations — language using must, shall, required to, or obligated to — and explicitly excludes explanatory text, preambles, and definitions.",
        "Sentences containing more than one distinct action are split into separate atomic obligations, while obligations that are logically inseparable are kept together.",
        "Obligations are grouped into between four and eight requirement clusters by topic, with two to six obligations per cluster, rather than grouped by paragraph order.",
        "An explicit deduplication pass removes obligations that share the same core action and subject within a group, keeping only one.",
        "Output remains in the document's detected source language, English or Arabic, with no translation applied.",
        "This stage runs at a temperature of 0.1, with a budget of up to 16,000 tokens.",
    ],
    accent=ORANGE,
)

bullets_slide(
    "UC-15 — Stage 2: Normalization and Classification",
    [
        "Stage 2 removes any exact duplicates that Stage 1 missed and splits any obligation that still contains more than one action.",
        "Each obligation is classified by type — Preventive, Detective, Governance, Reporting, or Documentation.",
        "Each obligation also receives a criticality of High, Medium, or Low, and an expected evidence type drawn from a fixed list: policy, procedure, system configuration, log, approval, report, contract, record, or other.",
        "An execution category is assigned to each obligation — Ongoing Control, One-Time Implementation, One-Off Reporting, Governance Approval, or Informational with No Action — and this single field determines whether Stage 3 runs at all.",
        "Each obligation receives a clarity score from one to five and a flag indicating whether it needs manual review because it cannot be made atomic.",
        "This stage runs at a temperature of 0.1, with a budget of up to 16,000 tokens.",
    ],
    accent=ORANGE,
)

bullets_slide(
    "UC-15 — Stage 3: Control Design",
    [
        "Stage 3 runs only for obligations classified as Ongoing Control in Stage 2.",
        "For each such obligation, the model designs one internal control, including a title, objective, a two-to-three sentence description, a realistic owning department, control type, execution type, frequency, control level, the evidence the control generates, three to five key steps, and the residual risk if the control fails.",
        "Obligations that are not Ongoing Control still pass through this stage, but with their control field left null — a control is never forced where one is not warranted.",
        "This stage runs at a temperature of 0.25, with a token budget that defaults to 8,000 since no override is set in the code.",
    ],
    accent=ORANGE,
)

bullets_slide(
    "UC-15 — Hallucination Control and Reliability",
    [
        "Every prompt carries an explicit anti-fabrication instruction: the system message states the model must never hallucinate, and Stage 1 specifically instructs it to extract only what is written, never invent obligations, and never paraphrase or interpret beyond the source text.",
        "Extraction runs at a temperature of 0.1, which sharply limits creative variation and produces near-deterministic output for the same input document.",
        "Stage 2 flags any obligation it cannot make fully atomic for manual review rather than guessing silently, and Stage 3 only designs a control for obligations Stage 2 has already classified as requiring one.",
        "Gap analysis applies the same grounding rule: the model is told to base every verdict only on the uploaded document's text and to quote a direct excerpt as evidence, never a paraphrase or an assumption.",
        "Requirement and control matching follow the same principle: the system prompt instructs the model to compare items precisely and without hallucination, and the model defaults to a verdict of \"new\" whenever a match is not clearly justified.",
        "There is no schema-constrained API output here, so resilience comes from prompt design plus a parser that strips formatting artifacts and falls back to an empty, clearly logged result rather than raising an error — a single malformed response is never retried into a fabricated success, and the pipeline simply moves on to the next document.",
    ],
    accent=ORANGE,
    notes="Grounded directly in the system prompts read from staged_LLM_Analyzer.py, gap_analyzer.py, and requirement_matcher.py — this consolidates every anti-hallucination and reliability mechanism actually present in the code.",
)

bullets_slide(
    "UC-15 — Matching, Gap Analysis, and Language Handling",
    [
        "Requirement matching calls the LLM once per extracted requirement to obtain a fully-matched, partially-matched, or new verdict against the bank's existing requirements, and then once more for each extracted control and each extracted KPI within that requirement — so matching cost scales with the number of requirements, not the number of documents.",
        "Each matching call is capped at 300 tokens, since the response is a short verdict and explanation rather than a generated document.",
        "Gap analysis compares an uploaded document against the extracted obligations. Uploads longer than 12,000 characters are split by paragraph and analyzed chunk by chunk, with results merged across chunks using a simple priority: covered outranks partial, and partial outranks missing.",
        "Language is detected once per document, and every stage prompt repeats the instruction that all output fields must remain in that language, with no translation applied — the engine classifies obligations in their original language rather than normalizing everything to English first.",
    ],
    accent=ORANGE,
)

table_slide(
    "UC-15 — Real Usage at Scale",
    ["Metric", "Value", "Note"],
    [
        ["Documents analyzed (3-stage pipeline)", "1,384 distinct docs", "Across roughly 4.4 months, January 28 to June 10, 2026"],
        ["Requirement rows produced", "2,713", "One row per requirement, not per obligation"],
        ["Activity pattern", "72% of all rows in a 3-day window (Apr 9–11)", "A backfill or reprocessing burst, not steady daily volume"],
        ["Documents that went through requirement matching", "24 of 1,384 (about 2%)", "Matching is barely exercised yet relative to analysis volume"],
        ["Controls / KPIs on file", "43 / 43 (31 AI-suggested each)", "—"],
    ],
    col_ratios=[3.4, 3.4, 4.6],
    accent=ORANGE,
    caption="These figures are pulled directly from compliance_analysis, sama_requirement_mapping, and the control and KPI tables.",
    notes="Pulled directly from compliance_analysis, sama_requirement_mapping, DEMO_CONTROL, DEMO_KPI.",
)

bullets_slide(
    "UC-15 — Lessons Learned",
    [
        "The model evaluation did not select the strongest model overall — Claude 3.5 Sonnet scored highest on analytical depth and forward-looking risk analysis. DeepSeek v3.2 was chosen instead because it delivered balanced performance across classification and structuring tasks while costing meaningfully less than Claude or GPT-4-class models, which matters for a pipeline that calls the model multiple times per document across thousands of documents.",
        "Because there is no schema-constrained output, the regular-expression-based JSON repair logic had to be designed deliberately rather than added as an afterthought; every stage already returns an empty result gracefully instead of crashing the document.",
        "Matching cost scales with the number of requirements rather than the number of documents — a single document with several requirement groups, each with its own controls and KPIs, can trigger a dozen or more additional LLM calls. That is likely why matching has, so far, only been run for twenty-four of the 1,384 analyzed documents.",
    ],
    accent=ORANGE,
)

# ───────────────────────── PART 3 — COSTING & CLOSE ─────────────────────────

table_slide(
    "Costing — Per-Document Economics",
    ["Regulator", "Avg. Cost per Document Analyzed", "Documents Analyzed to Date"],
    [
        ["CBB", "~$0.003", "9"],
        ["SBP", "~$0.008", "509"],
        ["SECP", "~$0.04", "195"],
        ["SAMA", "~$0.03", "671"],
        ["Total / Blended", "~$0.025 average", "1,384"],
    ],
    col_ratios=[1.8, 3.2, 3.0],
    accent=NAVY,
    max_row_height=Inches(0.85),
    caption="These are estimates derived from stored content length, not a metered invoice from OpenRouter. To project a monthly or annual figure, multiply the relevant cost per document by however many new regulations you expect to monitor in that period.",
    extra_bullets=[
        "Example: 50 new documents a month at SECP's rate (the most expensive, due to larger document sizes) would cost roughly $2 a month in LLM spend; the same 50 at SBP's rate would cost roughly $0.40 a month.",
        "The system has so far processed each regulator in concentrated catch-up batches rather than a steady monthly drip, so there is no observed steady-state publication rate to project from yet — the most reliable volume figure would come from each regulator's known historical publication frequency.",
        "This table covers UC-15's LLM analysis cost only. UC-2's other metered costs, 2Captcha and PDF.co, are real but are not currently logged in enough detail to compute an exact per-document figure — see the next slide.",
    ],
    notes="SAMA and SECP carry the highest per-document cost because their documents are genuinely larger, not because they're analyzed more often.",
)

bullets_slide(
    "Costing — What Isn't Instrumented",
    [
        "The number of CAPTCHA solves is not logged anywhere in the database. The processing log only records 347 download errors, a figure that combines network failures and CAPTCHA failures without distinguishing between them.",
        "Whether a given extraction used PDF.co or the free local OCR path is also not logged separately. The smart_extraction step recorded 202 successful extractions, but does not record which method actually succeeded.",
        "Exact LLM token counts are never captured from the OpenRouter response.",
        "For exact figures, the OpenRouter usage dashboard tied to the API key in the environment file reports exact token usage and spend, and the 2Captcha and PDF.co dashboards report historical credit and balance consumption — all three are a faster and more accurate source of truth than estimating from the database.",
    ],
    accent=NAVY,
)

bullets_slide(
    "Licensing and Infrastructure",
    [
        "A large share of the stack carries no license fee at all, including Scrapy, Selenium, Playwright, BeautifulSoup, PyMuPDF, pdfplumber, Tesseract OCR, FastAPI, SQLAlchemy, Streamlit, and deep-translator.",
        "Three services are paid and usage-based, with no subscription lock-in: OpenRouter for LLM tokens in UC-15, 2Captcha for CAPTCHA solving in the UC-2 crawlers, and PDF.co for format-preserving extraction in UC-2.",
        "The environment configuration currently points to a local SQL Server Express instance rather than a paid cloud database, so the real infrastructure cost today is effectively zero. That cost will only become real once the system is deployed to a hosted environment.",
    ],
    accent=NAVY,
)

bullets_slide(
    "Key Takeaways",
    [
        "UC-2 is the pipeline: four regulator-specific crawlers feed one unified flow that extracts, analyzes, stores, and serves regulatory content, with CBB alone carrying a complete version history.",
        "UC-15 is the engine inside that pipeline: a three-stage LLM process that extracts, classifies, and designs controls, combined with per-requirement matching and chunked gap analysis. The same engine serves both crawled and manually uploaded documents.",
        "Every major tool choice reflects a real, falsifiable decision: PDF.co was chosen for formatting fidelity after local OCR proved insufficient; DeepSeek v3.2 was chosen for balanced performance at meaningfully lower cost than Claude or GPT-4-class models; and each crawler framework was matched to its regulator's actual site behavior.",
        "Reliability is built in at every stage: explicit anti-hallucination instructions, low-temperature extraction, evidence-quoting rules, and graceful degradation on failure all work together so the system never guesses when it isn't sure.",
        "Real usage to date — 1,384 documents analyzed over four and a half months, at an estimated $0.025 average cost per document — is a small fraction of the 41,992 documents crawled. CBB's depth-based filter accounts for most of that difference.",
    ],
    accent=TEAL,
    notes="Close by reiterating the relationship: UC-2 is the pipeline, UC-15 is the brain inside it — and every number on the costing slides traces back to the actual database, not a generic estimate.",
)

slide = new_slide()
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
add_text(slide, Inches(1), Inches(3.3), Inches(11.33), Inches(1.0), "Thank you",
          size=40, bold=True, color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

OUTPUT_PATH = r"C:\Users\Hamdia\Desktop\UC-2_and_UC-15_Presentation_v7.pptx"
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
print("Final slide count:", len(prs.slides._sldIdLst))
